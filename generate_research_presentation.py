"""
科研工作汇报 PPT 生成脚本
基于曹祥龙医生团队的科研成果
使用 python-pptx 生成专业科研汇报演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # 设置标题样式
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # 设置副标题样式
    subtitle_frame = subtitle_shape.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(51, 51, 51)

def create_content_slide(prs, title, content_list, accent_color=RGBColor(0, 51, 102)):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = accent_color

    # 设置内容
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for item in content_list:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.level = 0

        # 检查是否是子项（以 - 开头）
        if item.strip().startswith('-'):
            p.level = 1
            p.font.size = Pt(18)

def create_two_column_slide(prs, title, left_content, right_content):
    """创建双栏对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 标题
    title_shape = slide.shapes.title
    title_shape.text = title

    # 左侧内容
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
    left_frame = left.text_frame
    left_frame.word_wrap = True

    p = left_frame.add_paragraph()
    p.text = left_content[0] if left_content else ""
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    for item in left_content[1:]:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_before = Pt(10)

    # 右侧内容
    right = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(5))
    right_frame = right.text_frame
    right_frame.word_wrap = True

    p = right_frame.add_paragraph()
    p.text = right_content[0] if right_content else ""
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    for item in right_content[1:]:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_before = Pt(10)

def create_highlight_slide(prs, title, highlights):
    """创建重点突出页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 标题
    title_shape = slide.shapes.title
    title_shape.text = title

    y_position = 1.5
    for highlight in highlights:
        # 添加圆角矩形背景
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5),
            Inches(y_position),
            Inches(9),
            Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(232, 244, 248)
        shape.line.color.rgb = RGBColor(0, 51, 102)

        # 添加文本
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = highlight
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)

        y_position += 1.0

def main():
    prs = Presentation()

    # 1. 标题页
    create_title_slide(
        prs,
        "老年胃肠肿瘤精准诊疗\n科研工作汇报",
        "曹祥龙医生团队\n北京医院 / 塔里木大学\n2026年1月"
    )

    # 2. 目录页
    create_content_slide(prs, "汇报内容", [
        "一、课题执行情况",
        "二、重点项目成果",
        "三、科研成果总结",
        "四、团队科研思路",
        "五、未来拓展规划",
        "六、总结与展望"
    ])

    # 3. 课题一
    create_content_slide(prs, "课题一：兵团重点领域科技攻关计划", [
        "项目名称：复合维生素调节肠道菌群TRG-5改善老年结直肠癌术后氧化应激反应的应用研究",
        "项目编号：2023AB018-131",
        "执行期限：2023.05.10 - 2026.05.09",
        "",
        "✅ 执行进展：",
        "- 2025年为关键攻坚年，进展顺利",
        "- 完成全部动物模型（POD3/5/16）样本采集",
        "- 完成临床队列样本的非靶向代谢组学检测",
        "- 完成16S rDNA测序",
        "",
        "🎯 关键成果：",
        "- 构建老年结直肠癌术后营养代谢多维数据库",
        "- 阐明'菌群-代谢'轴机制（TRG-5菌属→谷氨酰胺/丁酸通路）"
    ])

    # 4. 课题二
    create_content_slide(prs, "课题二：塔里木大学校长基金", [
        "项目名称：基于循环外泌体的胃癌脂质分子标志物筛选创新研究团队",
        "项目编号：TDZKCX202210 (2022ZD101)",
        "执行期限：2022.01.01 - 2024.12.31",
        "",
        "✅ 执行进展：",
        "- 项目已处于结题/执行报告阶段",
        "- 报告日期：2025年7月",
        "",
        "🎯 关键成果：",
        "- 建立外泌体分离纯化和脂质组学检测平台",
        "- 发现老年胃癌患者普遍存在脂质代谢紊乱",
        "- 超额完成科研产出：发表论文4篇（其中SCI 3篇）"
    ])

    # 5. 脂组学项目
    create_content_slide(prs, "脂组学项目 - 顺利结题", [
        "✅ 项目状态：已完成结题准备",
        "",
        "🔬 主要发现：",
        "- 建立外泌体分离纯化及质谱检测标准化流程",
        "- 揭示老年胃癌患者特有脂质代谢紊乱",
        "- 发现HexCer（棕榈酰脂肪酸）与胃癌转移潜在关联",
        "- 阐明ZDHHC4介导的Smad6棕榈酰化修饰机制",
        "",
        "📊 主要成果：",
        "- 学术：发表论文4篇（含3篇SCI），超额完成",
        "- 转化：开发基于机器学习的手术风险评估模型",
        "- 人才：培养硕士研究生2名"
    ])

    # 6. 维生素项目
    create_content_slide(prs, "维生素项目 - 按期执行", [
        "📌 当前阶段：数据挖掘与论文产出阶段",
        "",
        "📦 数据收集：",
        "- 动物实验：完成短肠综合征大鼠模型多时间点样本采集",
        "- 临床样本：完成双盲随机分组、干预及样本采集",
        "- 检测数据：完成代谢组学检测与生信分析",
        "",
        "🎯 预期成果：",
        "- 理论：明确复合维生素干预最佳剂量和时间窗",
        "- 应用：开发围手术期复合维生素制剂",
        "- 建立精准营养干预数学模型",
        "- 产出：发表高水平SCI论文1-2篇（已完成2篇）"
    ])

    # 7. 发表文章
    create_two_column_slide(
        prs,
        "已发表文章（2篇高水平SCI）",
        [
            "文章一：",
            "",
            "标题：Combined effects of depression and sedentary behavior on mortality risk",
            "期刊：BMC Geriatrics",
            "影响因子：~4.0 (JCR Q1)",
            "发表时间：2025.11.25"
        ],
        [
            "文章二：",
            "",
            "标题：Association between sarcopenic obesity and osteoarthritis",
            "期刊：Experimental Gerontology",
            "影响因子：Q1/Q2分区",
            "发表时间：2024.10.21"
        ]
    )

    # 8. 申请专利
    create_content_slide(prs, "申请专利（2项）", [
        "🏆 发明专利：",
        "- 名称：一种肠引流组件",
        "- 申请号：202410201886.X",
        "- 状态：申请已受理（2024.02.23）",
        "",
        "🏆 实用新型专利：",
        "- 名称：一种肠引流管",
        "- 申请号：202420343425.1",
        "- 状态：申请已受理（2024.02.26）",
        "",
        "💡 创新技术：'免还纳回肠造口'专利技术",
        "- 通过'T形引流+智能束紧'避免二次手术",
        "- 将单例费用降至传统手术的4%（约1000元 vs 2.5万元）"
    ])

    # 9. 科研思路体系
    create_highlight_slide(prs, "团队科研思路：全链条精准诊疗体系", [
        "🔍 上游（精准评估）：利用NHANES大数据挖掘，识别老年患者预后'隐形杀手'",
        "🔬 中游（机制解析）：代谢组学+肠道菌群测序，寻找早期诊断标志物和营养干预靶点",
        "🏥 下游（微创干预）：机器人手术+专利技术，实现极致精准+拒绝二次伤害"
    ])

    # 10. 技术路线
    create_content_slide(prs, "技术路线图", [
        "临床问题（大数据）",
        "↓",
        "科学假设",
        "↓",
        "技术方法（组学+测序）",
        "↓",
        "数据分析",
        "↓",
        "成果转化",
        "",
        "💡 核心创新：从基础发现到临床应用的完整闭环"
    ])

    # 11. 未来规划
    create_content_slide(prs, "2026年三大重点任务", [
        "📊 推数据（临床研究）：",
        "- 全力推进STARS-GC09多中心研究",
        "- 目标：完成300例患者入组",
        "- 产出：机器人手术安全性高质量循证证据",
        "",
        "🚀 推转化（技术推广）：",
        "- '免还纳回肠造口'技术申报医院新技术",
        "- 在医联体单位推广",
        "- 目标：吻合口漏率≤5%，费用降至传统手术4%",
        "",
        "📝 推结题（学术收官）：",
        "- 确保兵团科技攻关项目顺利结题",
        "- 发表'菌群-代谢轴'机制高分SCI论文",
        "- 申报相关发明专利"
    ])

    # 12. 总结
    create_highlight_slide(prs, "工作总结", [
        "✅ 2项兵团课题稳步推进，执行顺利",
        "✅ 脂组学项目成功结题，超额完成指标",
        "✅ 维生素项目按计划执行，数据完整",
        "✅ 发表高水平SCI论文2篇",
        "✅ 申请专利2项",
        "✅ 团队科研思路清晰，技术路线成熟"
    ])

    # 13. 创新亮点
    create_highlight_slide(prs, "创新亮点", [
        "🏆 建立'老年胃肠肿瘤全链条精准诊疗体系'",
        "🔬 发现HexCer脂质标志物和TRG-5菌群靶点",
        "💡 研发'免还纳回肠造口'专利技术",
        "🔄 形成从基础到临床的完整转化路径"
    ])

    # 14. 结束页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = "感谢聆听！"
    subtitle_shape.text = "敬请各位专家批评指正"

    # 保存文件
    output_file = "老年胃肠肿瘤科研工作汇报_自动生成版.pptx"
    prs.save(output_file)
    print(f"✅ PPT已生成：{output_file}")
    print(f"📊 共生成 {len(prs.slides)} 页幻灯片")

if __name__ == "__main__":
    main()
