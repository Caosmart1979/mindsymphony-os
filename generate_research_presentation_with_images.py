"""
科研工作汇报 PPT 生成脚本（带 AI 图片版本）
基于曹祥龙医生团队的科研成果
使用 python-pptx 生成专业科研汇报演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# 图片路径配置
IMAGE_DIR = "slide-images"
IMAGES = {
    "cover": os.path.join(IMAGE_DIR, "01-cover.png"),
    "research_system": os.path.join(IMAGE_DIR, "02-research-system.png"),
    "achievements": os.path.join(IMAGE_DIR, "03-achievements.png"),
    "patent_tech": os.path.join(IMAGE_DIR, "04-patent-tech.png"),
    "future_roadmap": os.path.join(IMAGE_DIR, "05-future-roadmap.png"),
    "lipidomics": os.path.join(IMAGE_DIR, "06-lipidomics-tech.png"),
}

def add_image_to_slide(slide, image_path, left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=None):
    """向幻灯片添加图片"""
    if os.path.exists(image_path):
        if height:
            slide.shapes.add_picture(image_path, left, top, width, height)
        else:
            slide.shapes.add_picture(image_path, left, top, width=width)
        return True
    else:
        print(f"⚠️ 图片不存在：{image_path}")
        return False

def create_title_slide_with_image(prs, title, subtitle, image_path):
    """创建带图片的标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 添加封面图片（全屏）
    if os.path.exists(image_path):
        add_image_to_slide(slide, image_path, Inches(0), Inches(0), Inches(10), Inches(7.5))
    else:
        # 如果图片不存在，使用传统文字标题
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[0]
        title_shape.text = title
        subtitle_shape.text = subtitle

def create_content_slide_with_image(prs, title, content_list, image_path=None, accent_color=RGBColor(0, 51, 102)):
    """创建带图片的内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 如果有图片，添加图片
    if image_path and os.path.exists(image_path):
        # 图片放在上半部分
        add_image_to_slide(slide, image_path, Inches(0.5), Inches(1.3), Inches(9), Inches(3.5))
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = accent_color

        # 内容放在下半部分
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for item in content_list[:5]:  # 限制显示前5条，避免溢出
            p = content_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.space_before = Pt(4)

            if item.strip().startswith('-'):
                p.level = 1
                p.font.size = Pt(14)
    else:
        # 无图片时使用标准布局
        create_content_slide_standard(slide, title, content_list, accent_color)

def create_content_slide_standard(slide, title, content_list, accent_color=RGBColor(0, 51, 102)):
    """标准内容页样式"""
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = accent_color

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for item in content_list:
        p = content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.space_before = Pt(8)

        if item.strip().startswith('-'):
            p.level = 1
            p.font.size = Pt(18)

def create_two_column_slide(prs, title, left_content, right_content):
    """创建双栏对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # 左侧内容
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
    left_frame = left.text_frame
    left_frame.word_wrap = True

    p = left_frame.add_paragraph()
    p.text = left_content[0] if left_content else ""
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    for item in left_content[1:]:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_before = Pt(10)

    # 右侧内容
    right = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(5))
    right_frame = right.text_frame
    right_frame.word_wrap = True

    p = right_frame.add_paragraph()
    p.text = right_content[0] if right_content else ""
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    for item in right_content[1:]:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_before = Pt(10)

def create_highlight_slide_with_image(prs, title, highlights, image_path=None):
    """创建带图片的重点突出页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # 如果有图片，添加图片
    if image_path and os.path.exists(image_path):
        add_image_to_slide(slide, image_path, Inches(0.5), Inches(1.3), Inches(9), Inches(3))
        y_position = 4.5
    else:
        y_position = 1.5

    # 添加重点内容
    for highlight in highlights:
        # 添加圆角矩形背景
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5),
            Inches(y_position),
            Inches(9),
            Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(232, 244, 248)
        shape.line.color.rgb = RGBColor(0, 51, 102)

        # 添加文本
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = highlight
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)

        y_position += 0.9

def create_image_slide(prs, title, image_path):
    """创建纯图片展示页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # 图片
    add_image_to_slide(slide, image_path, Inches(1), Inches(1.3), Inches(8), Inches(5.5))

def main():
    prs = Presentation()

    # 检查图片文件是否存在
    missing_images = []
    for key, path in IMAGES.items():
        if not os.path.exists(path):
            missing_images.append(f"{key}: {path}")

    if missing_images:
        print("⚠️ 警告：以下图片文件不存在：")
        for img in missing_images:
            print(f"  - {img}")
        print("\n将生成无图片版本的 PPT...\n")
    else:
        print("✅ 所有图片文件已找到，开始生成带图片的 PPT...\n")

    # 1. 封面页（带图片）
    if os.path.exists(IMAGES["cover"]):
        create_title_slide_with_image(
            prs,
            "老年胃肠肿瘤精准诊疗\n科研工作汇报",
            "曹祥龙医生团队 | 北京医院/塔里木大学",
            IMAGES["cover"]
        )
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = "老年胃肠肿瘤精准诊疗\n科研工作汇报"
        subtitle_shape.text = "曹祥龙医生团队\n北京医院 / 塔里木大学\n2026年1月"

    # 2. 目录页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "汇报内容"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    contents = [
        "一、课题执行情况",
        "二、重点项目成果",
        "三、科研成果总结",
        "四、团队科研思路",
        "五、未来拓展规划",
        "六、总结与展望"
    ]
    for item in contents:
        p = content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.space_before = Pt(12)

    # 3. 团队科研体系（带图片）
    create_image_slide(prs, "团队科研思路：全链条精准诊疗体系", IMAGES["research_system"])

    # 4. 课题一
    create_content_slide_standard(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "课题一：兵团重点领域科技攻关计划",
        [
            "项目名称：复合维生素调节肠道菌群TRG-5改善老年结直肠癌术后氧化应激反应的应用研究",
            f"项目编号：2023AB018-131",
            f"执行期限：2023.05.10 - 2026.05.09",
            "",
            "✅ 执行进展：",
            "- 2025年为关键攻坚年，进展顺利",
            "- 完成全部动物模型（POD3/5/16）样本采集",
            "- 完成临床队列样本的非靶向代谢组学检测",
            "- 完成16S rDNA测序",
            "",
            "🎯 关键成果：",
            "- 构建老年结直肠癌术后营养代谢多维数据库",
            "- 阐明'菌群-代谢'轴机制（TRG-5菌属→谷氨酰胺/丁酸通路）"
        ]
    )

    # 5. 课题二
    create_content_slide_standard(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "课题二：塔里木大学校长基金",
        [
            "项目名称：基于循环外泌体的胃癌脂质分子标志物筛选创新研究团队",
            "项目编号：TDZKCX202210 (2022ZD101)",
            "执行期限：2022.01.01 - 2024.12.31",
            "",
            "✅ 执行进展：",
            "- 项目已处于结题/执行报告阶段",
            "- 报告日期：2025年7月",
            "",
            "🎯 关键成果：",
            "- 建立外泌体分离纯化和脂质组学检测平台",
            "- 发现老年胃癌患者普遍存在脂质代谢紊乱",
            "- 超额完成科研产出：发表论文4篇（其中SCI 3篇）"
        ]
    )

    # 6. 脂组学项目（带图片）
    create_content_slide_with_image(
        prs,
        "脂组学项目 - 顺利结题",
        [
            "✅ 项目状态：已完成结题准备",
            "🔬 主要发现：HexCer棕榈酰脂肪酸（ZDHHC4机制）、TRG-5菌属、谷氨酰胺/丁酸代谢通路",
            "📊 主要成果：发表论文4篇（含3篇SCI），超额完成"
        ],
        IMAGES["lipidomics"]
    )

    # 7. 维生素项目
    create_content_slide_standard(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "维生素项目 - 按期执行",
        [
            "📌 当前阶段：数据挖掘与论文产出阶段",
            "",
            "📦 数据收集：",
            "- 动物实验：完成短肠综合征大鼠模型多时间点样本采集",
            "- 临床样本：完成双盲随机分组、干预及样本采集",
            "- 检测数据：完成代谢组学检测与生信分析",
            "",
            "🎯 预期成果：",
            "- 理论：明确复合维生素干预最佳剂量和时间窗",
            "- 应用：开发围手术期复合维生素制剂",
            "- 产出：发表高水平SCI论文1-2篇（已完成2篇）"
        ]
    )

    # 8. 发表文章
    create_two_column_slide(
        prs,
        "已发表文章（2篇高水平SCI）",
        [
            "文章一：",
            "",
            "标题：Combined effects of depression and sedentary behavior on mortality risk",
            "期刊：BMC Geriatrics",
            "影响因子：~4.0 (JCR Q1)",
            "发表时间：2025.11.25"
        ],
        [
            "文章二：",
            "",
            "标题：Association between sarcopenic obesity and osteoarthritis",
            "期刊：Experimental Gerontology",
            "影响因子：Q1/Q2分区",
            "发表时间：2024.10.21"
        ]
    )

    # 9. 成果总结（带图片）
    create_content_slide_with_image(
        prs,
        "学术成果总结",
        [
            "📚 论文：2篇高水平SCI（BMC Geriatrics、Experimental Gerontology）",
            "🏆 专利：2项专利申请（1项发明+1项实用新型）",
            "💡 创新：建立'老年胃肠肿瘤全链条精准诊疗体系'"
        ],
        IMAGES["achievements"]
    )

    # 10. 专利技术详解（带图片）
    create_content_slide_with_image(
        prs,
        "专利技术：免还纳回肠造口技术",
        [
            "💡 创新技术：'免还纳回肠造口'专利技术",
            "- 通过'T形引流+智能束紧'避免二次手术",
            "- 将单例费用降至传统手术的4%（约1000元 vs 2.5万元）",
            "🏆 专利申请：202410201886.X（发明）、202420343425.1（实用新型）"
        ],
        IMAGES["patent_tech"]
    )

    # 11. 未来规划（带图片）
    create_content_slide_with_image(
        prs,
        "2026年三大重点任务",
        [
            "📊 推数据：STARS-GC09多中心研究（300例）",
            "🚀 推转化：免还纳技术推广（吻合口漏率≤5%）",
            "📝 推结题：兵团课题结题（高分SCI）"
        ],
        IMAGES["future_roadmap"]
    )

    # 12. 总结
    create_highlight_slide_with_image(prs, "工作总结", [
        "✅ 2项兵团课题稳步推进，执行顺利",
        "✅ 脂组学项目成功结题，超额完成指标",
        "✅ 维生素项目按计划执行，数据完整",
        "✅ 发表高水平SCI论文2篇",
        "✅ 申请专利2项",
        "✅ 团队科研思路清晰，技术路线成熟"
    ])

    # 13. 创新亮点
    create_highlight_slide_with_image(prs, "创新亮点", [
        "🏆 建立'老年胃肠肿瘤全链条精准诊疗体系'",
        "🔬 发现HexCer脂质标志物和TRG-5菌群靶点",
        "💡 研发'免还纳回肠造口'专利技术",
        "🔄 形成从基础到临床的完整转化路径"
    ])

    # 14. 结束页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = "感谢聆听！"
    subtitle_shape.text = "敬请各位专家批评指正"

    # 保存文件
    output_file = "老年胃肠肿瘤科研工作汇报_AI图片版.pptx"
    prs.save(output_file)

    print(f"✅ PPT已生成：{output_file}")
    print(f"📊 共生成 {len(prs.slides)} 页幻灯片")
    print(f"🎨 包含 {len([i for i in IMAGES.values() if os.path.exists(i)])} 张 AI 图片")

if __name__ == "__main__":
    main()
